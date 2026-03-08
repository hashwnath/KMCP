"""Process uploaded documentation files into markdown-like text."""

from __future__ import annotations

import io
from pathlib import Path

import fitz
from docx import Document as DocxDocument
from pptx import Presentation


def _read_pdf(file_bytes: bytes) -> str:
    with fitz.open(stream=file_bytes, filetype="pdf") as pdf:
        return "\n\n".join(page.get_text("text") for page in pdf)


def _read_docx(file_bytes: bytes) -> str:
    document = DocxDocument(io.BytesIO(file_bytes))
    return "\n".join(paragraph.text for paragraph in document.paragraphs if paragraph.text)


def _read_pptx(file_bytes: bytes) -> str:
    presentation = Presentation(io.BytesIO(file_bytes))
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)


def process_file(file_bytes: bytes, filename: str) -> tuple[str, dict]:
    """Convert uploaded files into text suitable for chunking."""
    suffix = Path(filename).suffix.lower()
    if suffix == ".pdf":
        content = _read_pdf(file_bytes)
    elif suffix == ".docx":
        content = _read_docx(file_bytes)
    elif suffix == ".pptx":
        content = _read_pptx(file_bytes)
    else:
        content = file_bytes.decode("utf-8", errors="ignore")

    return content.strip(), {"title": Path(filename).stem, "filename": filename}